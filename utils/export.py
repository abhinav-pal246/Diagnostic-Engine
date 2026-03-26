import re
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Pt


def export_pdf(query: str, content: str) -> str:
    pdf = FPDF()
    pdf.add_page()

    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, "McKinsey Rapid Diagnostic", ln=True)

    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 10, query, ln=True)
    pdf.ln(5)

    pdf.set_font("Helvetica", size=10)
    clean = re.sub(r"\*\*(.*?)\*\*", r"\1", content)
    for line in clean.split("\n"):
        pdf.multi_cell(0, 7, line)

    path = "/tmp/diagnostic.pdf"
    pdf.output(path)
    return path


def export_pptx(query: str, content: str) -> str:
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "McKinsey Rapid Diagnostic"
    slide.placeholders[1].text = query

    # One slide per section
    titles = ["Situation", "Complication", "Resolution", "Key Risks"]
    for title in titles:
        pattern = rf"\*\*{title}\*\*(.*?)(?=\*\*|$)"
        match = re.search(pattern, content, re.DOTALL)
        if match:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            tf = slide.placeholders[1].text_frame
            tf.text = match.group(1).strip()
            tf.paragraphs[0].font.size = Pt(14)

    path = "/tmp/diagnostic.pptx"
    prs.save(path)
    return path